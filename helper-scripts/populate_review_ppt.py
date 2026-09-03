from pathlib import Path

from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Pt


ROOT = Path(__file__).resolve().parents[1]
SOURCE = ROOT / "Review1-phase3_ppt.pptx"
OUTPUT = ROOT / "Review1-phase3_ppt_completed.pptx"


def replace_text(shape, text, size=20):
    frame = shape.text_frame
    frame.clear()
    for index, line in enumerate(text.split("\n")):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = line
        paragraph.font.size = Pt(size)
        paragraph.font.name = "Aptos"
        paragraph.space_after = Pt(8)


presentation = Presentation(SOURCE)

replace_text(
    presentation.slides[0].shapes[0],
    "UE23AM441A  -  Capstone Project Phase - 3\n\nExplainable Multi-Agent ESG Risk Analysis System\nProject Progress Review #1",
    24,
)
replace_text(
    presentation.slides[0].shapes[1],
    "Project Title : Explainable Multi-Agent ESG Risk Analysis System\nProject ID : UE23AM441A\nProject Guide : [Add guide name]\nProject Team : [Add team member names]",
    18,
)

slide_content = {
    1: (
        "Outline",
        "Abstract and scope\nPhase 2 summary and improvements\nOverall architecture\nTasks, modules, and technologies\nIndividual contribution\nDemonstration and testing\nResults from 50-company run\nReferences",
    ),
    2: (
        "Abstract and Scope",
        "Develop an explainable seven-agent ESG analysis pipeline for Indian companies.\n\nScope:\n- Combine BRSR disclosures, AQI/location data, ESG news, and stock data.\n- Produce company-level ESG scores, quality gates, contradiction checks, and explanations.\n- Evaluate a configured cohort of 50 companies for FY 2024-25 where comparable data is available.",
    ),
    3: (
        "Summary of Work Done in Capstone Project Phase - 2",
        "Phase 2 established the V1 seven-agent ESG workflow for a five-company dataset.\n\nPhase 3 improvements:\n- Expanded the configured cohort from 5 to 50 companies.\n- Added validated BRSR ingestion and FY-aware source selection.\n- Added environmental risk, news sentiment, stock correlation, greenwashing checks, and quality gates.\n- Added output lineage, review flags, and an interactive web dashboard.",
    ),
    4: (
        "Architecture",
        "Input data\n  BRSR PDFs | AQI and location data | ESG news | Stock prices\n                         |\nSeven-agent pipeline\n  Extractor -> Environmental risk -> News sentiment -> Stock correlation\n  -> Greenwashing detector -> Master scorer -> Explainable AI\n                         |\nValidated CSV outputs + web dashboard + review flags",
    ),
    5: (
        "List of Tasks/Modules",
        "1. BRSR extraction and validation\n2. Environmental and location risk analysis\n3. ESG news filtering and sentiment classification\n4. Stock performance and ESG correlation\n5. Greenwashing and contradiction detection\n6. Master ESG scoring and quality gating\n7. Explainable AI report generation\n\nTechnologies: Python, Pandas, FastAPI, HTML/CSS/JavaScript, CSV/JSON, BRSR PDFs, AQI and market datasets.",
    ),
    6: (
        "Individual Contribution",
        "Contribution can be assigned and signed off by the team:\n\nAgent 1: BRSR extraction and validation\nAgent 2: Environmental risk features\nAgent 3: News sentiment pipeline\nAgent 4: Stock correlation analysis\nAgent 5: Greenwashing detection\nAgent 6: Master scoring and quality gate\nAgent 7: Explanations, reporting, and dashboard integration\n\nAdd member names, lines of code, and hours before submission.",
    ),
    7: (
        "Demonstration, Testing, and Results",
        "Demonstrated outputs from the Phase 3 run:\n- 50 company score records and 900 BRSR metric records\n- Average master ESG score: 6.17/10 for scored records\n- 677 ESG news sentiment records and 91 stock records\n- 83 environmental/location risk records\n- 4 companies passed the quality gate; 46 received low-quality warnings\n- Average BRSR completeness: 40.22%\n- Cross-validation flagged 2 high- and 6 medium-severity contradictions\n\nTesting: configuration count, PDF validation, data-quality reports, cross-validation, and API/dashboard checks.",
    ),
    8: (
        "References",
        "[1] Securities and Exchange Board of India, Business Responsibility and Sustainability Reporting (BRSR) framework.\n\n[2] Project datasets: BRSR reports, Air Quality Data in India (2015-2024), Real-Time AQI India, ESG News Dataset, and Sensex/Nifty historical data.\n\n[3] Python Software Foundation, Python documentation.\n\n[4] FastAPI, Pandas, and scikit-learn documentation.\n\n[5] Project implementation and validation reports, Phase 3, 2026.",
    ),
}

for slide_index, (title, body) in slide_content.items():
    slide = presentation.slides[slide_index]
    replace_text(slide.shapes[0], body, 16 if slide_index in (4, 6, 7, 8) else 18)
    replace_text(slide.shapes[1], title, 24)

for slide in list(presentation.slides)[1:9]:
    replace_text(slide.shapes[3], "Explainable Multi-Agent ESG Risk Analysis System", 10)
    replace_text(slide.shapes[4], "[Team names]", 10)

presentation.save(OUTPUT)
print(OUTPUT)